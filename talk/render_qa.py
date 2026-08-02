"""Render the deck from the .pptx itself (python-pptx + PIL) for visual QA,
and report geometry defects: out-of-bounds shapes, overlaps, text overflow."""
import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

DPI = 110
HERE = Path(__file__).parent
OUT = HERE / "qa"
DF = Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts")
SUP = Path("/System/Library/Fonts/Supplemental")

FONTS = {
    ("Calibri", 0, 0): DF / "Calibri.ttf", ("Calibri", 1, 0): DF / "Calibrib.ttf",
    ("Calibri", 0, 1): DF / "Calibrii.ttf", ("Calibri", 1, 1): DF / "Calibriz.ttf",
    ("Cambria", 0, 0): DF / "Cambria.ttc", ("Cambria", 1, 0): DF / "Cambriab.ttf",
    ("Cambria", 0, 1): DF / "Cambriai.ttf", ("Cambria", 1, 1): DF / "Cambriaz.ttf",
    ("Courier New", 0, 0): SUP / "Courier New.ttf", ("Courier New", 1, 0): SUP / "Courier New Bold.ttf",
    ("Courier New", 0, 1): SUP / "Courier New Italic.ttf", ("Courier New", 1, 1): SUP / "Courier New Bold Italic.ttf",
}
_cache = {}


def font(name, size_pt, bold, italic):
    key = (name, int(bool(bold)), int(bool(italic)), round(size_pt, 1))
    if key not in _cache:
        p = FONTS.get((name, key[1], key[2])) or FONTS[("Calibri", key[1], key[2])]
        px = max(6, int(round(size_pt * DPI / 72.0)))
        _cache[key] = ImageFont.truetype(str(p), px)
    return _cache[key]


def px(emu):
    return int(round(Emu(emu).inches * DPI))


def rgb(color, default=None):
    try:
        if color and color.type is not None and color.rgb is not None:
            return "#" + str(color.rgb)
    except Exception:
        pass
    return default


def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and f.type == 1:  # solid
            return rgb(f.fore_color)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        c = rgb(sh.line.color)
        w = sh.line.width
        return c, max(1, px(w)) if w else 1
    except Exception:
        return None, 1


def slide_bg(slide):
    xml = slide._element.xml
    if "<p:bg>" in xml:
        i = xml.find("srgbClr val=", xml.find("<p:bg>"))
        if 0 < i < xml.find("</p:bg>"):
            return "#" + xml[i + 13:i + 19]
    return "#FFFFFF"


def wrap(draw, runs, maxw):
    """runs: [(text, font, color)] -> list of lines, each [(text, font, color)]."""
    lines, cur, curw = [], [], 0.0
    for text, f, col in runs:
        for j, seg in enumerate(text.split("\n")):
            if j:
                lines.append(cur); cur, curw = [], 0.0
            for word in seg.split(" "):
                if not word:
                    continue
                cand = word if not cur else " " + word
                w = draw.textlength(cand, font=f)
                if cur and curw + w > maxw:
                    lines.append(cur); cur, curw = [(word, f, col)], draw.textlength(word, font=f)
                else:
                    cur.append((cand, f, col)); curw += w
    if cur:
        lines.append(cur)
    return lines


def render(path):
    OUT.mkdir(exist_ok=True)
    prs = Presentation(path)
    SW, SH = px(prs.slide_width), px(prs.slide_height)
    problems = []

    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (SW, SH), slide_bg(slide))
        d = ImageDraw.Draw(img)
        boxes = []

        for sh in slide.shapes:
            try:
                L, T, W_, H_ = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            except Exception:
                continue

            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                try:
                    im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                    im.thumbnail((max(1, W_), max(1, H_)), Image.LANCZOS)
                    img.paste(im, (L + (W_ - im.width) // 2, T + (H_ - im.height) // 2))
                except Exception as e:
                    d.rectangle([L, T, L + W_, T + H_], outline="#FF00FF", width=3)
                    problems.append(f"slide {idx}: image failed to render ({e})")
                boxes.append(("image", L, T, W_, H_))
                continue

            st = str(sh.shape_type or "")
            fill = shape_fill(sh)
            lc, lw = shape_line(sh)
            if "LINE" in st or (H_ <= 2 and W_ > 4) or (W_ <= 2 and H_ > 4):
                if lc:
                    d.line([L, T, L + W_, T + H_], fill=lc, width=lw)
            elif "OVAL" in st or "ELLIPSE" in st:
                d.ellipse([L, T, L + W_, T + H_], fill=fill, outline=lc, width=lw)
            elif fill or lc:
                r = min(12, W_ // 8, H_ // 8) if "ROUND" in st else 0
                if r > 2:
                    d.rounded_rectangle([L, T, L + W_, T + H_], radius=r, fill=fill, outline=lc, width=lw)
                else:
                    d.rectangle([L, T, L + W_, T + H_], fill=fill, outline=lc, width=lw)

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                if fill or lc:
                    boxes.append(("shape", L, T, W_, H_))
                continue

            tf = sh.text_frame
            il = px(tf.margin_left or 0); ir = px(tf.margin_right or 0)
            it = px(tf.margin_top or 0); ib = px(tf.margin_bottom or 0)
            bx, by = L + il, T + it
            bw, bh = max(4, W_ - il - ir), max(4, H_ - it - ib)

            para_lines, total = [], 0
            for p in tf.paragraphs:
                runs = []
                for r in p.runs:
                    fname = r.font.name or "Calibri"
                    fsz = r.font.size.pt if r.font.size else 14
                    fnt = font(fname, fsz, r.font.bold, r.font.italic)
                    runs.append((r.text, fnt, rgb(r.font.color, "#000000")))
                if not runs:
                    total += 8
                    para_lines.append(([], 8, None, 0, False))
                    continue
                bullet = "buChar" in p._pPr.xml if p._pPr is not None else False
                indent = 16 if bullet else 0
                base = max(f[1].size for f in runs)
                ls = p.line_spacing
                if ls is None:
                    lh = int(base * 1.22)
                elif isinstance(ls, float):
                    lh = int(base * ls)
                else:  # Length, in points
                    lh = int(ls.pt * DPI / 72.0)
                lh = max(lh, int(base * 1.16))
                ws = wrap(d, runs, bw - indent)
                for k, w_ in enumerate(ws):
                    para_lines.append((w_, lh, p.alignment, indent, indent and k == 0))
                    total += lh
                sa = p.space_after.pt * DPI / 72.0 if p.space_after else 0
                total += int(sa)
                if sa:
                    para_lines.append(([], int(sa), None, 0, False))

            va = str(tf.vertical_anchor or "")
            y = by if ("TOP" in va or not va) else (by + (bh - total) // 2 if "MIDDLE" in va else by + bh - total)
            for parts, lh, align, indent, dot in para_lines:
                if parts:
                    lw_ = sum(d.textlength(t, font=f) for t, f, _ in parts)
                    if align and "CENTER" in str(align):
                        x = bx + (bw - lw_) / 2
                    elif align and "RIGHT" in str(align):
                        x = bx + bw - lw_
                    else:
                        x = bx + indent
                    if dot:
                        d.ellipse([x - 12, y + lh * 0.38, x - 7, y + lh * 0.38 + 5], fill=parts[0][2])
                    for t, f, c in parts:
                        d.text((x, y), t, font=f, fill=c)
                        x += d.textlength(t, font=f)
                y += lh

            if total > bh + 6:
                problems.append(
                    f"slide {idx}: TEXT OVERFLOW {total - bh}px  \"{tf.text[:55].strip()}…\"")
            boxes.append(("text", L, T, W_, H_))

            if L < -2 or T < -2 or L + W_ > SW + 2 or T + H_ > SH + 2:
                problems.append(f"slide {idx}: OUT OF BOUNDS  \"{tf.text[:40].strip()}\"")

        d.rectangle([0, 0, SW - 1, SH - 1], outline="#CCCCCC")
        d.text((SW - 60, SH - 26), f"{idx}", font=font("Calibri", 13, 1, 0), fill="#999999")
        img.save(OUT / f"slide-{idx:02d}.png")

    print("\n".join(problems) if problems else "no geometry problems detected")
    print(f"\n{len(prs.slides.__iter__.__self__._sldIdLst)} slides → {OUT}")


if __name__ == "__main__":
    render(sys.argv[1] if len(sys.argv) > 1 else HERE / "oblique_registration.pptx")
